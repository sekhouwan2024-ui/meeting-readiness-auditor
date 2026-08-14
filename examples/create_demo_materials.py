#!/usr/bin/env python3
"""Create synthetic PPTX and XLSX inputs for a safe local demo."""
from __future__ import annotations

import argparse
from pathlib import Path

from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt


def build_workbook(path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Operating Summary"
    sheet.append(["Metric", "April", "May", "Target"])
    sheet.append(["Revenue", 100000, 112000, 115000])
    sheet.append(["Gross margin", 0.46, 0.42, 0.45])
    sheet.append(["Paid D1 retention", 0.31, 0.26, 0.30])
    sheet.append(["Ad spend", 20000, 28000, 25000])
    for cell in sheet[3] + sheet[4]:
        if isinstance(cell.value, float):
            cell.number_format = "0.0%"
    workbook.save(path)


def build_presentation(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(8.8), Inches(0.8))
    title.text_frame.paragraphs[0].text = "May Operating Review"
    title.text_frame.paragraphs[0].font.size = Pt(30)
    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(3.8))
    body.text_frame.text = (
        "Revenue grew 12% month over month.\n"
        "Gross margin remained healthy at 46%.\n"
        "Paid D1 retention was broadly stable.\n"
        "Recommendation: increase ad spend by another 20%."
    )
    presentation.save(path)


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", required=True)
    args = parser.parse_args()
    output = Path(args.out).expanduser().resolve()
    output.mkdir(parents=True, exist_ok=True)
    build_workbook(output / "operating-data.xlsx")
    build_presentation(output / "operating-review.pptx")
    print(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
