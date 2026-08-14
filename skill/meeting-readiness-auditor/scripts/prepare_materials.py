#!/usr/bin/env python3
"""Prepare a stable evidence workspace for meeting-readiness auditing.

The script performs deterministic file extraction and writes:
- workspace_manifest.json: paths and schema metadata
- inventory.json: concise file inventory
- analysis_packet.json: compact model-facing packet
- evidence.sqlite: complete searchable evidence store
- previews/: rendered PPT slide images when system tools are available

It does not make business judgments and never emits a monolithic source bundle.
"""
from __future__ import annotations

import argparse
import csv
import hashlib
import json
import re
import shutil
import sqlite3
import subprocess
import tempfile
from datetime import date, datetime
from pathlib import Path
from typing import Any, Iterable

SCHEMA_VERSION = "2.0"
SUPPORTED = {".pptx", ".xlsx", ".xlsm", ".csv", ".docx", ".pdf"}
NUMERIC_LINE_RE = re.compile(
    r"(?i)(?:[-+]?\d[\d,]*(?:\.\d+)?\s*(?:%|pp|x|倍|万|万元|元|人|单|个|天|小时|亿|千)?|同比|环比|增长|下降|提升|减少|达到|完成率|转化率|留存率|利润率|毛利率|ROI|ROAS|CAC|GMV)"
)


def json_value(value: Any) -> Any:
    if isinstance(value, (datetime, date)):
        return value.isoformat()
    if isinstance(value, bytes):
        return value.decode("utf-8", errors="replace")
    if isinstance(value, (int, float, bool)) or value is None:
        return value
    return str(value)


def compact_text(value: Any, limit: int = 1200) -> str:
    text = "" if value is None else str(value)
    text = re.sub(r"\s+", " ", text).strip()
    return text if len(text) <= limit else text[: limit - 1] + "…"


def collect_paths(inputs: list[str]) -> list[Path]:
    found: list[Path] = []
    for raw in inputs:
        path = Path(raw).expanduser().resolve()
        if path.is_dir():
            found.extend(item for item in path.rglob("*") if item.is_file())
        elif path.is_file():
            found.append(path)
    unique = sorted({item for item in found if item.suffix.lower() in SUPPORTED}, key=lambda p: str(p).lower())
    return unique


def deterministic_file_id(path: Path, index: int) -> str:
    digest = hashlib.sha1(str(path).encode("utf-8")).hexdigest()[:8]
    return f"F{index:03d}-{digest}"


def ensure_schema(connection: sqlite3.Connection) -> None:
    connection.executescript(
        """
        PRAGMA journal_mode=WAL;
        PRAGMA synchronous=NORMAL;
        CREATE TABLE metadata (key TEXT PRIMARY KEY, value TEXT NOT NULL);
        CREATE TABLE files (
            file_id TEXT PRIMARY KEY,
            file_name TEXT NOT NULL,
            file_type TEXT NOT NULL,
            absolute_path TEXT NOT NULL,
            role_hint TEXT,
            extraction_status TEXT NOT NULL,
            warning TEXT
        );
        CREATE TABLE ppt_blocks (
            source_id TEXT PRIMARY KEY,
            file_id TEXT NOT NULL,
            page INTEGER NOT NULL,
            block_type TEXT NOT NULL,
            object_name TEXT,
            text_content TEXT,
            json_payload TEXT,
            left_pos INTEGER,
            top_pos INTEGER,
            width INTEGER,
            height INTEGER
        );
        CREATE INDEX idx_ppt_file_page ON ppt_blocks(file_id, page);
        CREATE INDEX idx_ppt_text ON ppt_blocks(text_content);
        CREATE TABLE spreadsheet_cells (
            source_id TEXT PRIMARY KEY,
            file_id TEXT NOT NULL,
            sheet_name TEXT NOT NULL,
            address TEXT NOT NULL,
            row_index INTEGER NOT NULL,
            column_index INTEGER NOT NULL,
            value_text TEXT,
            numeric_value REAL,
            formula TEXT,
            number_format TEXT
        );
        CREATE INDEX idx_cell_file_sheet ON spreadsheet_cells(file_id, sheet_name, row_index, column_index);
        CREATE INDEX idx_cell_text ON spreadsheet_cells(value_text);
        CREATE TABLE document_blocks (
            source_id TEXT PRIMARY KEY,
            file_id TEXT NOT NULL,
            page_or_section TEXT,
            block_type TEXT NOT NULL,
            text_content TEXT,
            json_payload TEXT
        );
        CREATE INDEX idx_doc_file ON document_blocks(file_id, page_or_section);
        CREATE INDEX idx_doc_text ON document_blocks(text_content);
        """
    )
    connection.execute("INSERT INTO metadata(key, value) VALUES (?, ?)", ("schema_version", SCHEMA_VERSION))


def render_pptx(path: Path, preview_dir: Path) -> dict[str, Any]:
    preview_dir.mkdir(parents=True, exist_ok=True)
    soffice = shutil.which("libreoffice") or shutil.which("soffice")
    pdftoppm = shutil.which("pdftoppm")
    if not soffice or not pdftoppm:
        return {"rendered": False, "reason": "preview tools unavailable", "images": []}
    try:
        with tempfile.TemporaryDirectory(prefix="meeting-audit-") as temp_dir:
            result = subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", temp_dir, str(path)],
                capture_output=True,
                text=True,
                timeout=120,
            )
            pdf_path = Path(temp_dir) / f"{path.stem}.pdf"
            if result.returncode != 0 or not pdf_path.exists():
                return {"rendered": False, "reason": compact_text(result.stderr or result.stdout, 300), "images": []}
            prefix = preview_dir / "slide"
            result = subprocess.run(
                [pdftoppm, "-png", "-r", "120", str(pdf_path), str(prefix)],
                capture_output=True,
                text=True,
                timeout=120,
            )
            if result.returncode != 0:
                return {"rendered": False, "reason": compact_text(result.stderr or result.stdout, 300), "images": []}
        images = sorted(preview_dir.glob("slide-*.png"))
        return {"rendered": bool(images), "reason": None, "images": [str(item) for item in images]}
    except Exception as exc:
        return {"rendered": False, "reason": f"{type(exc).__name__}: {exc}", "images": []}


def insert_pptx(
    connection: sqlite3.Connection,
    path: Path,
    file_id: str,
    preview_root: Path,
    render: bool,
) -> tuple[dict[str, Any], list[dict[str, Any]], list[str]]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for PPTX input") from exc

    presentation = Presentation(str(path))
    slide_summaries: list[dict[str, Any]] = []
    visual_queue: list[dict[str, Any]] = []
    warnings: list[str] = []

    preview = {"rendered": False, "images": [], "reason": "render disabled"}
    if render:
        preview = render_pptx(path, preview_root / file_id)
        if not preview.get("rendered") and preview.get("reason"):
            warnings.append(str(preview["reason"]))

    for page, slide in enumerate(presentation.slides, start=1):
        text_items: list[dict[str, Any]] = []
        tables_count = 0
        charts_summary: list[dict[str, Any]] = []
        image_count = 0
        numeric_claims: list[str] = []
        for shape_index, shape in enumerate(slide.shapes, start=1):
            source_prefix = f"{file_id}:P{page}:O{shape_index}"
            shape_name = getattr(shape, "name", f"Object {shape_index}")
            if getattr(shape, "shape_type", None) == 13:  # picture
                image_count += 1
            if getattr(shape, "has_text_frame", False):
                raw_text = "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()
                if raw_text:
                    connection.execute(
                        """INSERT INTO ppt_blocks(source_id,file_id,page,block_type,object_name,text_content,json_payload,left_pos,top_pos,width,height)
                           VALUES(?,?,?,?,?,?,?,?,?,?,?)""",
                        (
                            f"{source_prefix}:text",
                            file_id,
                            page,
                            "text",
                            shape_name,
                            raw_text,
                            None,
                            int(shape.left),
                            int(shape.top),
                            int(shape.width),
                            int(shape.height),
                        ),
                    )
                    text_items.append(
                        {
                            "object_name": shape_name,
                            "text": compact_text(raw_text, 800),
                            "top": int(shape.top),
                            "left": int(shape.left),
                        }
                    )
                    for line in raw_text.splitlines():
                        line = compact_text(line, 400)
                        if line and NUMERIC_LINE_RE.search(line) and line not in numeric_claims:
                            numeric_claims.append(line)
            if getattr(shape, "has_table", False):
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                connection.execute(
                    """INSERT INTO ppt_blocks(source_id,file_id,page,block_type,object_name,text_content,json_payload,left_pos,top_pos,width,height)
                       VALUES(?,?,?,?,?,?,?,?,?,?,?)""",
                    (
                        f"{source_prefix}:table",
                        file_id,
                        page,
                        "table",
                        shape_name,
                        "\n".join(" | ".join(row) for row in rows),
                        json.dumps({"rows": rows}, ensure_ascii=False),
                        int(shape.left),
                        int(shape.top),
                        int(shape.width),
                        int(shape.height),
                    ),
                )
                tables_count += 1
            if getattr(shape, "has_chart", False):
                chart = shape.chart
                chart_info: dict[str, Any] = {
                    "object_name": shape_name,
                    "chart_type": str(chart.chart_type),
                    "title": None,
                    "series": [],
                    "value_axis": None,
                    "category_axis": None,
                }
                try:
                    if chart.has_title:
                        chart_info["title"] = chart.chart_title.text_frame.text
                except Exception:
                    pass
                try:
                    categories = [str(category.label) for category in chart.plots[0].categories]
                except Exception:
                    categories = []
                all_numeric_values: list[float] = []
                for series in chart.series:
                    try:
                        raw_values = [json_value(value) for value in series.values]
                    except Exception:
                        raw_values = []
                    for value in raw_values:
                        if isinstance(value, (int, float)):
                            all_numeric_values.append(float(value))
                    chart_info["series"].append(
                        {"name": str(series.name), "categories": categories, "values": raw_values}
                    )
                try:
                    axis = chart.value_axis
                    chart_info["value_axis"] = {
                        "minimum_scale": json_value(axis.minimum_scale),
                        "maximum_scale": json_value(axis.maximum_scale),
                        "major_unit": json_value(axis.major_unit),
                        "has_title": bool(axis.has_title),
                        "title": axis.axis_title.text_frame.text if axis.has_title else None,
                    }
                except Exception:
                    pass
                try:
                    axis = chart.category_axis
                    chart_info["category_axis"] = {
                        "has_title": bool(axis.has_title),
                        "title": axis.axis_title.text_frame.text if axis.has_title else None,
                    }
                except Exception:
                    pass
                connection.execute(
                    """INSERT INTO ppt_blocks(source_id,file_id,page,block_type,object_name,text_content,json_payload,left_pos,top_pos,width,height)
                       VALUES(?,?,?,?,?,?,?,?,?,?,?)""",
                    (
                        f"{source_prefix}:chart",
                        file_id,
                        page,
                        "chart",
                        shape_name,
                        compact_text(chart_info.get("title") or shape_name),
                        json.dumps(chart_info, ensure_ascii=False),
                        int(shape.left),
                        int(shape.top),
                        int(shape.width),
                        int(shape.height),
                    ),
                )
                axis_info = chart_info.get("value_axis") or {}
                axis_min = axis_info.get("minimum_scale")
                axis_max = axis_info.get("maximum_scale")
                truncated_candidate = (
                    isinstance(axis_min, (int, float))
                    and float(axis_min) > 0
                    and all_numeric_values
                    and min(all_numeric_values) >= 0
                )
                charts_summary.append(
                    {
                        "object_name": shape_name,
                        "title": chart_info.get("title"),
                        "chart_type": chart_info.get("chart_type"),
                        "series_count": len(chart_info["series"]),
                        "value_axis": chart_info.get("value_axis"),
                        "truncated_axis_candidate": bool(truncated_candidate),
                    }
                )
                if truncated_candidate:
                    visual_queue.append(
                        {
                            "file_id": file_id,
                            "file_name": path.name,
                            "page": page,
                            "object_name": shape_name,
                            "reason": "value_axis_min_above_zero",
                            "axis_min": axis_min,
                            "axis_max": axis_max,
                            "action": "Check whether the truncated axis exaggerates the visual change.",
                        }
                    )
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            notes = ""
        if notes:
            connection.execute(
                "INSERT INTO ppt_blocks(source_id,file_id,page,block_type,object_name,text_content,json_payload) VALUES(?,?,?,?,?,?,?)",
                (f"{file_id}:P{page}:notes", file_id, page, "notes", "notes", notes, None),
            )
        sorted_texts = sorted(text_items, key=lambda item: (item["top"], item["left"]))
        title = sorted_texts[0]["text"].split("\n")[0] if sorted_texts else f"第{page}页"
        preview_path = None
        images = preview.get("images", [])
        if page - 1 < len(images):
            preview_path = images[page - 1]
        if image_count and not charts_summary:
            visual_queue.append(
                {
                    "file_id": file_id,
                    "file_name": path.name,
                    "page": page,
                    "reason": "image_based_content",
                    "action": "Inspect the rendered slide because charts or tables may be images.",
                }
            )
        slide_summaries.append(
            {
                "page": page,
                "title": compact_text(title, 180),
                "key_texts": [item["text"] for item in sorted_texts[:8]],
                "numeric_claims": numeric_claims[:16],
                "table_count": tables_count,
                "chart_count": len(charts_summary),
                "image_count": image_count,
                "charts": charts_summary,
                "preview_path": preview_path,
            }
        )
    return (
        {
            "file_id": file_id,
            "file_name": path.name,
            "file_type": "pptx",
            "slide_count": len(slide_summaries),
            "slides": slide_summaries,
            "preview": preview,
        },
        visual_queue,
        warnings,
    )


def spreadsheet_row_values(ws: Any, formula_ws: Any, row_index: int, max_columns: int = 20) -> list[Any]:
    values: list[Any] = []
    for column_index in range(1, min(ws.max_column, max_columns) + 1):
        value = ws.cell(row=row_index, column=column_index).value
        formula_value = formula_ws.cell(row=row_index, column=column_index).value
        if value is None and isinstance(formula_value, str) and formula_value.startswith("="):
            value = formula_value
        values.append(json_value(value))
    while values and values[-1] is None:
        values.pop()
    return values


def insert_xlsx(
    connection: sqlite3.Connection,
    path: Path,
    file_id: str,
    max_cells_per_sheet: int,
) -> tuple[dict[str, Any], list[str]]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("openpyxl is required for XLSX input") from exc

    workbook_values = load_workbook(path, data_only=True, read_only=False)
    workbook_formulas = load_workbook(path, data_only=False, read_only=False)
    sheet_summaries: list[dict[str, Any]] = []
    warnings: list[str] = []
    for worksheet in workbook_values.worksheets:
        formula_sheet = workbook_formulas[worksheet.title]
        inserted = 0
        total_non_empty = 0
        formula_count = 0
        truncated = False
        non_empty_rows: list[int] = []
        for row in worksheet.iter_rows():
            row_has_value = False
            for cell in row:
                formula_raw = formula_sheet[cell.coordinate].value
                if cell.value is None and formula_raw is None:
                    continue
                row_has_value = True
                total_non_empty += 1
                if inserted >= max_cells_per_sheet:
                    truncated = True
                    continue
                formula = formula_raw if isinstance(formula_raw, str) and formula_raw.startswith("=") else None
                if formula:
                    formula_count += 1
                numeric_value = float(cell.value) if isinstance(cell.value, (int, float)) and not isinstance(cell.value, bool) else None
                value_text = None if cell.value is None else str(json_value(cell.value))
                source_id = f"{file_id}:{worksheet.title}:{cell.coordinate}"
                connection.execute(
                    """INSERT INTO spreadsheet_cells(source_id,file_id,sheet_name,address,row_index,column_index,value_text,numeric_value,formula,number_format)
                       VALUES(?,?,?,?,?,?,?,?,?,?)""",
                    (
                        source_id,
                        file_id,
                        worksheet.title,
                        cell.coordinate,
                        cell.row,
                        cell.column,
                        value_text,
                        numeric_value,
                        formula,
                        cell.number_format,
                    ),
                )
                inserted += 1
            if row_has_value:
                non_empty_rows.append(row[0].row)
        preview_rows = []
        metric_rows = []
        for row_index in non_empty_rows[:30]:
            values = spreadsheet_row_values(worksheet, formula_sheet, row_index)
            if values:
                preview_rows.append({"row": row_index, "values": values})
        for row_index in non_empty_rows[:300]:
            values = spreadsheet_row_values(worksheet, formula_sheet, row_index)
            has_label = any(isinstance(value, str) and value.strip() for value in values)
            has_number = any(isinstance(value, (int, float)) and not isinstance(value, bool) for value in values)
            has_formula = any(isinstance(value, str) and value.startswith("=") for value in values)
            if has_label and (has_number or has_formula):
                metric_rows.append({"row": row_index, "values": values})
            if len(metric_rows) >= 40:
                break
        if truncated:
            warnings.append(f"{worksheet.title}: non-empty cells exceeded {max_cells_per_sheet}; evidence store is truncated")
        sheet_summaries.append(
            {
                "sheet_name": worksheet.title,
                "dimensions": {"max_row": worksheet.max_row, "max_column": worksheet.max_column},
                "non_empty_cells": total_non_empty,
                "stored_cells": inserted,
                "formula_count": formula_count,
                "truncated": truncated,
                "preview_rows": preview_rows,
                "metric_rows": metric_rows,
            }
        )
    workbook_values.close()
    workbook_formulas.close()
    return (
        {
            "file_id": file_id,
            "file_name": path.name,
            "file_type": "xlsx",
            "sheet_count": len(sheet_summaries),
            "sheets": sheet_summaries,
        },
        warnings,
    )


def insert_csv(connection: sqlite3.Connection, path: Path, file_id: str, max_rows: int) -> tuple[dict[str, Any], list[str]]:
    rows_preview: list[dict[str, Any]] = []
    rows_count = 0
    truncated = False
    with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as handle:
        reader = csv.reader(handle)
        for row_index, row in enumerate(reader, start=1):
            if row_index > max_rows:
                truncated = True
                break
            rows_count += 1
            if row_index <= 30:
                rows_preview.append({"row": row_index, "values": row[:20]})
            for column_index, value in enumerate(row, start=1):
                if value == "":
                    continue
                address = f"R{row_index}C{column_index}"
                connection.execute(
                    """INSERT INTO spreadsheet_cells(source_id,file_id,sheet_name,address,row_index,column_index,value_text,numeric_value,formula,number_format)
                       VALUES(?,?,?,?,?,?,?,?,?,?)""",
                    (
                        f"{file_id}:CSV:{address}",
                        file_id,
                        "CSV",
                        address,
                        row_index,
                        column_index,
                        value,
                        None,
                        None,
                        None,
                    ),
                )
    return (
        {
            "file_id": file_id,
            "file_name": path.name,
            "file_type": "csv",
            "row_count": rows_count,
            "truncated": truncated,
            "preview_rows": rows_preview,
        },
        ["CSV row limit reached"] if truncated else [],
    )


def insert_docx(connection: sqlite3.Connection, path: Path, file_id: str) -> tuple[dict[str, Any], list[str]]:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is required for DOCX input") from exc
    document = Document(str(path))
    paragraph_count = 0
    table_count = 0
    preview: list[str] = []
    for index, paragraph in enumerate(document.paragraphs, start=1):
        text = paragraph.text.strip()
        if not text:
            continue
        paragraph_count += 1
        if len(preview) < 20:
            preview.append(compact_text(text, 500))
        connection.execute(
            "INSERT INTO document_blocks(source_id,file_id,page_or_section,block_type,text_content,json_payload) VALUES(?,?,?,?,?,?)",
            (f"{file_id}:P{index}", file_id, f"paragraph {index}", "paragraph", text, None),
        )
    for table_index, table in enumerate(document.tables, start=1):
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        table_count += 1
        connection.execute(
            "INSERT INTO document_blocks(source_id,file_id,page_or_section,block_type,text_content,json_payload) VALUES(?,?,?,?,?,?)",
            (
                f"{file_id}:T{table_index}",
                file_id,
                f"table {table_index}",
                "table",
                "\n".join(" | ".join(row) for row in rows),
                json.dumps({"rows": rows}, ensure_ascii=False),
            ),
        )
    return (
        {
            "file_id": file_id,
            "file_name": path.name,
            "file_type": "docx",
            "paragraph_count": paragraph_count,
            "table_count": table_count,
            "preview_paragraphs": preview,
        },
        [],
    )


def insert_pdf(connection: sqlite3.Connection, path: Path, file_id: str) -> tuple[dict[str, Any], list[str]]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("pypdf is required for PDF input") from exc
    reader = PdfReader(str(path))
    previews: list[dict[str, Any]] = []
    warnings: list[str] = []
    for page_index, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception as exc:
            text = ""
            warnings.append(f"page {page_index}: {type(exc).__name__}: {exc}")
        connection.execute(
            "INSERT INTO document_blocks(source_id,file_id,page_or_section,block_type,text_content,json_payload) VALUES(?,?,?,?,?,?)",
            (f"{file_id}:P{page_index}", file_id, f"page {page_index}", "page", text, None),
        )
        if page_index <= 20:
            previews.append({"page": page_index, "text": compact_text(text, 700)})
    return (
        {
            "file_id": file_id,
            "file_name": path.name,
            "file_type": "pdf",
            "page_count": len(reader.pages),
            "preview_pages": previews,
        },
        warnings,
    )


def role_hint(path: Path) -> str:
    lower = path.name.lower()
    if path.suffix.lower() == ".pptx":
        return "report_candidate"
    if any(term in lower for term in ["口径", "定义", "说明"]):
        return "definition_candidate"
    if path.suffix.lower() in {".xlsx", ".xlsm", ".csv"}:
        return "data_candidate"
    return "supporting_material"


def trim_packet(packet: dict[str, Any], char_budget: int) -> dict[str, Any]:
    serialized = json.dumps(packet, ensure_ascii=False)
    if len(serialized) <= char_budget:
        return packet
    packet["packet_truncated"] = True
    for file_summary in packet.get("file_summaries", []):
        if file_summary.get("file_type") == "xlsx":
            for sheet in file_summary.get("sheets", []):
                sheet["preview_rows"] = sheet.get("preview_rows", [])[:10]
                sheet["metric_rows"] = sheet.get("metric_rows", [])[:15]
        elif file_summary.get("file_type") == "pptx":
            for slide in file_summary.get("slides", []):
                slide["key_texts"] = slide.get("key_texts", [])[:4]
                slide["numeric_claims"] = slide.get("numeric_claims", [])[:8]
    serialized = json.dumps(packet, ensure_ascii=False)
    if len(serialized) <= char_budget:
        return packet
    for file_summary in packet.get("file_summaries", []):
        if file_summary.get("file_type") == "xlsx":
            for sheet in file_summary.get("sheets", []):
                sheet.pop("preview_rows", None)
        elif file_summary.get("file_type") == "pptx":
            for slide in file_summary.get("slides", []):
                slide.pop("key_texts", None)
    return packet


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("inputs", nargs="+", help="Files or directories to inspect")
    parser.add_argument("--out", required=True, help="Workspace output directory")
    parser.add_argument("--max-cells-per-sheet", type=int, default=200000)
    parser.add_argument("--max-csv-rows", type=int, default=100000)
    parser.add_argument("--packet-char-budget", type=int, default=140000)
    parser.add_argument("--no-render", action="store_true")
    args = parser.parse_args()

    output_dir = Path(args.out).expanduser().resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    paths = collect_paths(args.inputs)
    if not paths:
        print(json.dumps({"ok": False, "error": "No supported input files found"}, ensure_ascii=False))
        return 2

    database_path = output_dir / "evidence.sqlite"
    if database_path.exists():
        database_path.unlink()
    connection = sqlite3.connect(database_path)
    ensure_schema(connection)

    file_summaries: list[dict[str, Any]] = []
    file_inventory: list[dict[str, Any]] = []
    visual_review_queue: list[dict[str, Any]] = []
    global_warnings: list[dict[str, Any]] = []

    for index, path in enumerate(paths, start=1):
        file_id = deterministic_file_id(path, index)
        extension = path.suffix.lower()
        file_type = extension.lstrip(".")
        hint = role_hint(path)
        status = "ok"
        warning_text = None
        summary: dict[str, Any] | None = None
        file_warnings: list[str] = []
        try:
            if extension == ".pptx":
                summary, visual_items, file_warnings = insert_pptx(
                    connection,
                    path,
                    file_id,
                    output_dir / "previews",
                    not args.no_render,
                )
                visual_review_queue.extend(visual_items)
            elif extension in {".xlsx", ".xlsm"}:
                summary, file_warnings = insert_xlsx(connection, path, file_id, args.max_cells_per_sheet)
            elif extension == ".csv":
                summary, file_warnings = insert_csv(connection, path, file_id, args.max_csv_rows)
            elif extension == ".docx":
                summary, file_warnings = insert_docx(connection, path, file_id)
            elif extension == ".pdf":
                summary, file_warnings = insert_pdf(connection, path, file_id)
        except Exception as exc:
            status = "failed"
            warning_text = f"{type(exc).__name__}: {exc}"
            global_warnings.append({"file_id": file_id, "file_name": path.name, "message": warning_text})
        if file_warnings:
            global_warnings.extend(
                {"file_id": file_id, "file_name": path.name, "message": message} for message in file_warnings
            )
            warning_text = "; ".join(file_warnings)
        connection.execute(
            "INSERT INTO files(file_id,file_name,file_type,absolute_path,role_hint,extraction_status,warning) VALUES(?,?,?,?,?,?,?)",
            (file_id, path.name, file_type, str(path), hint, status, warning_text),
        )
        inventory_item = {
            "file_id": file_id,
            "file_name": path.name,
            "file_type": file_type,
            "role_hint": hint,
            "extraction_status": status,
            "warnings": file_warnings or ([warning_text] if warning_text else []),
        }
        file_inventory.append(inventory_item)
        if summary is not None:
            summary["role_hint"] = hint
            file_summaries.append(summary)

    connection.commit()
    counts = {
        "files": connection.execute("SELECT COUNT(*) FROM files WHERE extraction_status='ok'").fetchone()[0],
        "ppt_blocks": connection.execute("SELECT COUNT(*) FROM ppt_blocks").fetchone()[0],
        "spreadsheet_cells": connection.execute("SELECT COUNT(*) FROM spreadsheet_cells").fetchone()[0],
        "document_blocks": connection.execute("SELECT COUNT(*) FROM document_blocks").fetchone()[0],
    }
    connection.close()

    workspace_id = hashlib.sha1("|".join(item["file_id"] for item in file_inventory).encode("utf-8")).hexdigest()[:12]
    inventory = {
        "schema_version": SCHEMA_VERSION,
        "workspace_id": workspace_id,
        "generated_at": datetime.now().isoformat(timespec="seconds"),
        "files": file_inventory,
        "counts": counts,
        "warnings": global_warnings,
    }
    analysis_packet = {
        "schema_version": SCHEMA_VERSION,
        "workspace_id": workspace_id,
        "purpose": "Compact model-facing evidence packet. Query evidence.sqlite for details; do not load the database or a monolithic source dump into context.",
        "file_summaries": file_summaries,
        "visual_review_queue": visual_review_queue,
        "query_examples": [
            "python scripts/query_evidence.py <workspace> --query GMV --limit 30",
            "python scripts/query_evidence.py <workspace> --query 投放费用 --type xlsx_cell --limit 30",
            "python scripts/query_evidence.py <workspace> --file-id F001-xxxx --page 5 --limit 50",
        ],
        "packet_truncated": False,
    }
    analysis_packet = trim_packet(analysis_packet, args.packet_char_budget)

    inventory_path = output_dir / "inventory.json"
    packet_path = output_dir / "analysis_packet.json"
    manifest_path = output_dir / "workspace_manifest.json"
    inventory_path.write_text(json.dumps(inventory, ensure_ascii=False, indent=2), encoding="utf-8")
    packet_path.write_text(json.dumps(analysis_packet, ensure_ascii=False, indent=2), encoding="utf-8")
    manifest = {
        "schema_version": SCHEMA_VERSION,
        "workspace_id": workspace_id,
        "inventory": str(inventory_path),
        "analysis_packet": str(packet_path),
        "evidence_database": str(database_path),
        "previews_directory": str(output_dir / "previews"),
        "counts": counts,
        "warnings_count": len(global_warnings),
    }
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps({"ok": counts["files"] > 0, "manifest": str(manifest_path), "counts": counts, "warnings": len(global_warnings)}, ensure_ascii=False))
    return 0 if counts["files"] > 0 else 2


if __name__ == "__main__":
    raise SystemExit(main())
